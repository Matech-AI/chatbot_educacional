"""
Document converter: normalizes all ingestion formats to Markdown before RAG indexing.

Supported formats:
  - .md, .txt       → passthrough (already text)
  - .pdf            → markitdown → Markdown  (fallback: PyPDF2)
  - .docx, .doc     → markitdown → Markdown
  - .pptx, .ppt     → markitdown → Markdown  (fallback: python-pptx slide text)
  - .srt            → custom SRT parser → Markdown transcript
  - others          → markitdown best-effort, skipped on failure

Converted files are cached in <cache_dir>/<stem>_<hash12>.md so re-indexing is
instant when the source file has not changed.
"""

import hashlib
import logging
import re
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# markitdown availability
# ---------------------------------------------------------------------------
try:
    from markitdown import MarkItDown
    _MARKITDOWN_AVAILABLE = True
    logger.info("✅ markitdown disponível — conversão normalizada habilitada")
except ImportError:
    _MARKITDOWN_AVAILABLE = False
    logger.warning(
        "⚠️ markitdown não instalado. Instale com: pip install markitdown\n"
        "   A indexação usará os loaders legados para PDF/MD e ignorará DOCX/SRT."
    )


# ---------------------------------------------------------------------------
# Fallback helpers (used when markitdown is absent)
# ---------------------------------------------------------------------------

def _pdf_fallback(path: Path) -> Optional[str]:
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n\n".join(f"## Página {i+1}\n\n{text}" for i, text in enumerate(pages) if text.strip())
    except Exception:
        pass
    try:
        import PyPDF2
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            pages = [reader.pages[i].extract_text() or "" for i in range(len(reader.pages))]
        return "\n\n".join(f"## Página {i+1}\n\n{text}" for i, text in enumerate(pages) if text.strip())
    except Exception as e:
        logger.debug(f"PDF fallback failed for {path.name}: {e}")
        return None


def _pptx_fallback(path: Path) -> Optional[str]:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides.append(f"## Slide {i}\n\n" + "\n\n".join(texts))
        return "\n\n".join(slides) if slides else None
    except Exception as e:
        logger.debug(f"PPTX fallback failed for {path.name}: {e}")
        return None


# ---------------------------------------------------------------------------
# SRT parser (no external dependency needed)
# ---------------------------------------------------------------------------

def _parse_srt(path: Path) -> str:
    """Convert SRT subtitle file to a clean Markdown transcript."""
    try:
        raw = path.read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        logger.warning(f"⚠️ Cannot read SRT {path.name}: {e}")
        return ""

    # Split into blocks separated by blank lines
    blocks = re.split(r"\n{2,}", raw.strip())
    lines: list[str] = [f"# Transcrição: {path.stem}\n"]

    for block in blocks:
        parts = block.strip().splitlines()
        # SRT block: index, timestamp line, text line(s)
        if len(parts) >= 3 and "-->" in parts[1]:
            text = " ".join(p.strip() for p in parts[2:] if p.strip())
        elif len(parts) == 2 and "-->" in parts[0]:
            text = parts[1].strip()
        elif len(parts) >= 2 and "-->" not in parts[0]:
            text = " ".join(p.strip() for p in parts if p.strip())
        else:
            continue
        # Strip HTML tags common in SRT (<i>, <b>, <font ...>)
        text = re.sub(r"<[^>]+>", "", text).strip()
        if text:
            lines.append(text)

    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Main converter class
# ---------------------------------------------------------------------------

# Extensions that pass through as-is (already text/Markdown)
_PASSTHROUGH_EXTS = {".md", ".txt"}

# Extensions handled by markitdown (or their fallbacks)
_CONVERTIBLE_EXTS = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".srt", ".html", ".htm"}


class DocumentConverter:
    """
    Converts documents of various formats into Markdown strings.

    Converted Markdown is cached to <cache_dir> using the source file's
    content hash so re-processing is a no-op when content hasn't changed.
    """

    def __init__(self, cache_dir: str = "data/converted"):
        self.cache_dir = Path(cache_dir)
        self.cache_dir.mkdir(parents=True, exist_ok=True)

        if _MARKITDOWN_AVAILABLE:
            try:
                self._md = MarkItDown()
            except Exception as e:
                logger.warning(f"⚠️ MarkItDown init failed: {e}")
                self._md = None
        else:
            self._md = None

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    @property
    def available(self) -> bool:
        return _MARKITDOWN_AVAILABLE and self._md is not None

    @staticmethod
    def is_supported(path: Path) -> bool:
        return path.suffix.lower() in (_PASSTHROUGH_EXTS | _CONVERTIBLE_EXTS)

    def convert(self, source_path: Path) -> Optional[str]:
        """
        Return Markdown text for *source_path*, or None if conversion fails.

        For passthrough extensions (.md, .txt) the file is read directly.
        For all other supported types the result is cached on disk.
        """
        ext = source_path.suffix.lower()

        # Passthrough — no conversion needed, no caching overhead
        if ext in _PASSTHROUGH_EXTS:
            try:
                return source_path.read_text(encoding="utf-8", errors="replace")
            except Exception as e:
                logger.warning(f"⚠️ Cannot read {source_path.name}: {e}")
                return None

        if ext not in _CONVERTIBLE_EXTS:
            return None

        # Check disk cache
        cached = self._cached_markdown(source_path)
        if cached is not None:
            return cached

        # Dispatch conversion
        if ext == ".srt":
            result = _parse_srt(source_path)
        elif self._md is not None:
            result = self._convert_with_markitdown(source_path)
            if result is None:
                result = self._fallback(source_path)
        else:
            result = self._fallback(source_path)

        # Write cache if we got something useful
        if result and len(result.strip()) > 50:
            try:
                cache_path = self._cache_path(source_path)
                cache_path.write_text(result, encoding="utf-8")
            except Exception as e:
                logger.debug(f"Cache write failed for {source_path.name}: {e}")

        return result or None

    # ------------------------------------------------------------------
    # Internal helpers
    # ------------------------------------------------------------------

    def _file_hash(self, path: Path) -> str:
        try:
            return hashlib.md5(path.read_bytes()).hexdigest()[:12]
        except Exception:
            return hashlib.md5(str(path).encode()).hexdigest()[:12]

    def _cache_path(self, source_path: Path) -> Path:
        h = self._file_hash(source_path)
        return self.cache_dir / f"{source_path.stem}_{h}.md"

    def _cached_markdown(self, source_path: Path) -> Optional[str]:
        try:
            cache_path = self._cache_path(source_path)
            if cache_path.exists():
                text = cache_path.read_text(encoding="utf-8")
                logger.debug(f"📦 Cache hit: {source_path.name}")
                return text
        except Exception:
            pass
        return None

    def _convert_with_markitdown(self, path: Path) -> Optional[str]:
        try:
            result = self._md.convert(str(path))
            text = result.text_content
            if text and len(text.strip()) > 30:
                logger.debug(f"✅ markitdown: {path.name} ({len(text)} chars)")
                return text
            return None
        except Exception as e:
            logger.debug(f"markitdown failed for {path.name}: {e}")
            return None

    def _fallback(self, path: Path) -> Optional[str]:
        """Format-specific fallbacks used when markitdown is unavailable or fails."""
        ext = path.suffix.lower()
        if ext == ".pdf":
            return _pdf_fallback(path)
        if ext in (".pptx", ".ppt"):
            return _pptx_fallback(path)
        # DOCX: try python-docx
        if ext in (".docx", ".doc"):
            try:
                import docx2txt
                text = docx2txt.process(str(path))
                return text if text and text.strip() else None
            except Exception:
                pass
            try:
                from docx import Document as DocxDocument
                doc = DocxDocument(str(path))
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                return "\n\n".join(paragraphs) or None
            except Exception as e:
                logger.debug(f"DOCX fallback failed for {path.name}: {e}")
                return None
        return None
